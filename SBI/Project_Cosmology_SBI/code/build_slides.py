"""build_slides.py — generate the presentation deck (python-pptx)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
FIG  = os.path.join(HERE, "..", "figures")
OUTP = os.path.join(HERE, "..", "Presentation_Cosmology_SBI.pptx")

NAVY=RGBColor(0x12,0x1A,0x3A); INK=RGBColor(0x1B,0x24,0x4A)
LIGHT=RGBColor(0xF4,0xF6,0xFB); WHITE=RGBColor(0xFF,0xFF,0xFF)
ICE=RGBColor(0xCA,0xDC,0xFC); GOLD=RGBColor(0xE8,0xA1,0x3A); TEAL=RGBColor(0x1C,0x72,0x93)
MUTE=RGBColor(0x5A,0x63,0x80)
HEAD="Cambria"; BODY="Calibri"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
def blank():
    s=prs.slides.add_slide(prs.slide_layouts[6]); return s
def bg(s,c):
    s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def tb(s,l,t,w,h,lines,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,wrap=True):
    b=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=b.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    for i,(txt,sz,bold,col) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=txt; f=r.font
        f.size=Pt(sz); f.bold=bold; f.color.rgb=col; f.name=BODY
    return b
def head(s,txt,col=INK):
    tb(s,0.6,0.42,12.1,1.0,[(txt,30,True,col)])
def pic(s,path,left,top,width):
    im=Image.open(path); ar=im.height/im.width
    w=Inches(width); h=Emu(int(w*ar))
    s.shapes.add_picture(path,Inches(left),Inches(top),width=w)
    return width, w/914400.0*ar  # inches w,h
def cap(s,txt,l,t,w):
    tb(s,l,t,w,0.4,[(txt,11,False,MUTE)],align=PP_ALIGN.CENTER)
def chip(s,l,t,txt):
    sp=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(l),Inches(t),Inches(0.5),Inches(0.5))
    sp.fill.solid(); sp.fill.fore_color.rgb=GOLD; sp.line.fill.background()
    tf=sp.text_frame; tf.word_wrap=False; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=txt; r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=NAVY; r.font.name=BODY

# ---------- S1 Title ----------
s=blank(); bg(s,NAVY)
tb(s,0.9,2.0,11.5,2.2,[
  ("Estimating Cosmological Parameters",40,True,WHITE),
  ("from the Matter Power Spectrum",40,True,WHITE),
  ("Amortized simulation-based inference with neural posterior estimation",18,False,ICE)])
tb(s,0.9,4.7,11.5,1.4,[
  ("Group members:  Md Saiful Islam  ·  Sumon Ahmed Masum  ·  Mohammed Rehan",16,False,ICE),
  ("Simulation-Based Inference (SoSe26) · TU Dortmund · Supervisor: Aayush",13,False,MUTE)])
tb(s,0.9,6.6,11.5,0.5,[("θ = { H₀ , Ωₘ , nₛ , Aₛ }   ←   P(k)",18,True,GOLD)])

# ---------- S2 Motivation ----------
s=blank(); bg(s,LIGHT); head(s,"Why this problem matters")
tb(s,0.6,1.5,6.0,5.2,[
 ("The Universe is not smooth — matter clumps into the cosmic web.",16,True,INK),
 ("The matter power spectrum P(k) measures how much structure exists at each scale; galaxy surveys (DESI, Euclid) measure it.",14,False,INK),
 ("",8,False,INK),
 ("The inverse problem:",15,True,TEAL),
 ("We can simulate parameters → P(k) (forward), but we want P(k) → parameters, with honest uncertainty (backward).",14,False,INK),
 ("",8,False,INK),
 ("Why SBI:",15,True,TEAL),
 ("• Realistic simulators have no tractable likelihood.",14,False,INK),
 ("• Amortized inference: train once, infer any new spectrum instantly.",14,False,INK),
 ("• Calibrated error bars → trustworthy cosmology (e.g. the Hubble tension).",14,False,INK)])
w,h=pic(s,os.path.join(FIG,"fig2_observation.png"),6.9,1.7,6.0); cap(s,"One noisy observation of P(k) (cosmic-variance noise)",6.9,1.7+h+0.05,6.0)

# ---------- S3 Parameters ----------
s=blank(); bg(s,LIGHT); head(s,"The four parameters shape P(k)")
w,h=pic(s,os.path.join(FIG,"fig1_simulator_effects.png"),0.5,1.5,7.4)
cap(s,"Each parameter leaves a distinct imprint on the spectrum",0.5,1.5+h+0.03,7.4)
tb(s,8.2,1.6,4.7,5.4,[
 ("Aₛ  amplitude",15,True,GOLD),("vertical scaling of P(k)",13,False,INK),("",6,False,INK),
 ("nₛ  spectral index",15,True,GOLD),("tilts / pivots the spectrum",13,False,INK),("",6,False,INK),
 ("Ωₘ  matter density",15,True,GOLD),("moves the turnover scale",13,False,INK),("",6,False,INK),
 ("H₀  Hubble constant",15,True,GOLD),("weak effect — enters only via shape → hard to constrain alone",13,False,INK)])

# ---------- S4 Method ----------
s=blank(); bg(s,LIGHT); head(s,"Method: amortized neural posterior estimation")
# soft translucent-style fills (light tints) with coloured borders
LNAVY=RGBColor(0xDD,0xE4,0xF3); LTEAL=RGBColor(0xD8,0xEC,0xF2); LGOLD=RGBColor(0xF9,0xEC,0xD3)
def box(l,t,w,h,txt,fill,border,fg,sz=14,bw=1.5):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    sp.line.color.rgb=border; sp.line.width=Pt(bw)
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=txt
    r.font.size=Pt(sz); r.font.bold=True; r.font.color.rgb=fg; r.font.name=BODY
def arrow(l,t,w):
    sp=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(l),Inches(t),Inches(w),Inches(0.5))
    sp.fill.solid(); sp.fill.fore_color.rgb=GOLD; sp.line.fill.background()
box(0.7,2.1,3.0,1.1,"Prior\nθ = {H₀,Ωₘ,nₛ,Aₛ}",LNAVY,NAVY,NAVY)
arrow(3.85,2.35,1.0)
box(5.0,2.1,3.4,1.1,"Simulator\nCAMB / BBKS + noise",LTEAL,TEAL,TEAL)
arrow(8.55,2.35,1.0)
box(9.7,2.1,3.0,1.1,"P(k)\nnoisy spectrum",LNAVY,NAVY,NAVY)
box(5.0,4.2,3.4,1.1,"Neural Posterior\nq(θ | P(k))",LGOLD,GOLD,INK,14,2.25)
arrow(8.55,4.45,1.0); box(9.7,4.2,3.0,1.1,"posterior\nover θ",LNAVY,NAVY,NAVY)
# down arrow text
tb(s,0.7,4.1,4.0,1.2,[("Train once on many simulations →",14,True,INK),
  ("inference for any new spectrum is one fast forward pass (amortized).",13,False,INK)])
tb(s,0.7,5.7,12,0.5,[("Built with:  CAMB / BBKS simulator  ·  Gaussian NPE (Python, NumPy, autograd)  ·  BayesFlow for production  ·  Go + HTML live demo",12.5,False,INK)])
tb(s,0.7,6.35,12,0.8,[("Likelihood-free: we only need to simulate, never to write p(data | θ).  Maps to course Ch10 (ABC) → Ch12 (NPE/BayesFlow).",13,False,MUTE)])

# ---------- S5 Simulator, priors, noise ----------
s=blank(); bg(s,LIGHT); head(s,"Simulator, priors & noise")
tb(s,0.6,1.5,6.2,5.3,[
 ("Forward model",15,True,TEAL),
 ("Linear P(k) at z=0 on 100 log-spaced k ∈ [10⁻³,1] h/Mpc.",13,False,INK),
 ("Prototype: BBKS+Sugiyama transfer fn (≈ CAMB);  production: CAMB.",13,False,INK),
 ("",6,False,INK),
 ("Priors (proper, uniform)",15,True,TEAL),
 ("H₀ ~ U(60, 80) km/s/Mpc",13,False,INK),
 ("Ωₘ ~ U(0.10, 0.50)",13,False,INK),
 ("nₛ ~ U(0.90, 1.05)",13,False,INK),
 ("ln(10¹⁰Aₛ) ~ U(2.5, 3.5)   (log-amplitude)",13,False,INK),
 ("fixed ω_b = 0.0224, ω_c derived",13,False,MUTE)])
tb(s,7.0,1.5,5.9,5.3,[
 ("Cosmic-variance noise",15,True,TEAL),
 ("σᵢ = P(kᵢ) √(2 / Nᵢ)",15,True,GOLD),
 ("Nᵢ = kᵢ² Δk V / (2π²),   V = 10⁹ (Mpc/h)³",13,False,INK),
 ("Few modes at low k → noisy;  many at high k → precise.",13,False,INK),
 ("",6,False,INK),
 ("Preprocessing",15,True,TEAL),
 ("log₁₀ P(k), then standardize (raw P spans orders of magnitude).",13,False,INK),
 ("",6,False,INK),
 ("Data: 12k train / 2k val / 3k test;  fresh noise each batch.",13,False,INK)])

# ---------- S6 Network & training ----------
s=blank(); bg(s,LIGHT); head(s,"Network & training")
tb(s,0.55,1.35,6.25,5.7,[
 ("Architecture (≈31k params)",14,True,TEAL),
 ("• Input: standardized log P(k) — 100 values",12.5,False,INK),
 ("• MLP: 100 → 128 → 128, tanh  (hidden layers = summary)",12.5,False,INK),
 ("• Output 14 = mean(4) + Cholesky(10) → 4-D Gaussian, full covariance",12.5,False,INK),
 ("• No dropout / weight decay — online noise acts as regularization",12.5,False,INK),
 ("",4,False,INK),
 ("Training",14,True,TEAL),
 ("• Loss: Gaussian negative log-likelihood (NLL)",12.5,False,INK),
 ("• Two-phase: MSE warm-up (2.5k) → NLL (7k steps)",12.5,False,INK),
 ("• Adam, batch 256, step LR decay (2e-3 → 1e-3)",12.5,False,INK),
 ("• Offline sims + fresh online noise per batch",12.5,False,INK),
 ("• ~20 s on CPU;  val tracks train → no overfit",12.5,False,INK),
 ("",4,False,INK),
 ("Tools: Python · NumPy · autograd",12.5,False,MUTE),
 ("Production → CNN summary + BayesFlow normalizing flow",12.5,True,GOLD)])
w,h=pic(s,os.path.join(FIG,"fig3_loss.png"),7.05,2.0,5.85); cap(s,"Training / validation convergence",7.05,2.0+h+0.03,5.85)
# network internals (forward pass) — bottom strip
tb(s,0.55,5.4,9,0.4,[("Network internals (forward pass)",13,True,TEAL)])
box(0.55,5.8,1.65,0.8,"P(k)\n100",LNAVY,NAVY,NAVY,12)
arrow(2.27,6.1,0.42)
box(2.78,5.8,1.8,0.8,"Dense 128\ntanh",LTEAL,TEAL,TEAL,12)
arrow(4.65,6.1,0.42)
box(5.15,5.8,1.8,0.8,"Dense 128\ntanh",LTEAL,TEAL,TEAL,12)
arrow(7.02,6.1,0.42)
box(7.5,5.8,2.4,0.8,"14 outputs\nmean + Cholesky",LGOLD,GOLD,INK,12)
arrow(9.97,6.1,0.42)
box(10.45,5.8,2.45,0.8,"4-D Gaussian\nposterior",LNAVY,NAVY,NAVY,12)

# ---------- S7 Calibration ----------
s=blank(); bg(s,LIGHT); head(s,"Calibration — is the posterior trustworthy?")
w,h=pic(s,os.path.join(FIG,"fig5_sbc.png"),0.5,1.4,12.3)
w2,h2=pic(s,os.path.join(FIG,"fig6_contraction.png"),0.5,4.25,4.5)
tb(s,5.6,4.45,7.3,2.6,[
 ("nₛ, Aₛ ranks flat → well calibrated.",14,True,INK),
 ("Ωₘ, H₀ mild ∪-shape → slight tail overconfidence (typical for a Gaussian posterior; a normalizing flow would improve it).",13,False,INK),
 ("z-score std ≈ 0.95, 95% coverage ≈ 0.96–1.00.",13,False,TEAL),
 ("Contraction high for Ωₘ/nₛ/Aₛ, ≈ 0 for H₀.",13,False,INK)])

# ---------- S8 Recovery ----------
s=blank(); bg(s,LIGHT); head(s,"Parameter recovery")
w,h=pic(s,os.path.join(FIG,"fig4_recovery.png"),0.5,1.55,12.0)
cap(s,"Posterior mean vs true (test set). r = correlation, rmse = error.",0.5,1.55+h+0.04,12.0)
tb(s,0.6,5.35,12.4,0.6,[("Ωₘ:  r = 0.97          nₛ:  r = 0.91          Aₛ:  r = 0.87          — strongly recovered.",17,True,TEAL)])
tb(s,0.6,6.05,12.4,0.6,[("H₀:  r = 0.30  — weakly constrained, as physics predicts (next slide).",16,True,GOLD)])

# ---------- S9 Fiducial inference + degeneracy ----------
s=blank(); bg(s,LIGHT); head(s,"Inference at the Planck-2018 cosmology")
w,h=pic(s,os.path.join(FIG,"fig7_fiducial_corner.png"),0.5,1.5,5.3)
cap(s,"Posterior corner; red = truth",0.5,1.5+h+0.02,5.3)
tb(s,6.4,1.6,6.5,5.4,[
 ("Truth is recovered:",15,True,TEAL),
 ("Ωₘ = 0.305 ± 0.033  (true 0.315)",14,False,INK),
 ("nₛ = 0.972 ± 0.014  (true 0.965)",14,False,INK),
 ("H₀ = 69.0 ± 6.4   (true 67.4) — broad",14,False,INK),
 ("",6,False,INK),
 ("The H₀–Ωₘ degeneracy",16,True,GOLD),
 ("Correlation = −0.96.  The linear spectrum's shape depends on Γ ≈ Ωₘ h, so H₀ and Ωₘ trade off along constant Γ.",14,False,INK),
 ("A broad, calibrated H₀ is the correct scientific result — not a failure.",14,True,TEAL)])

# ---------- S10 PPC + limitations ----------
s=blank(); bg(s,LIGHT); head(s,"Model check & limitations")
w,h=pic(s,os.path.join(FIG,"fig8_ppc.png"),0.5,1.6,6.0); cap(s,"Posterior predictive spectra bracket the observation",0.5,1.6+h+0.03,6.0)
tb(s,7.0,1.6,5.9,5.3,[
 ("Posterior predictive check passes",15,True,TEAL),
 ("Predicted spectra tightly bracket the observed P(k) at all scales.",14,False,INK),
 ("",8,False,INK),
 ("Honest limitations",15,True,GOLD),
 ("• Gaussian posterior → mild tail overconfidence (→ use a normalizing flow / BayesFlow).",14,False,INK),
 ("• BBKS analytic simulator, linear only → use CAMB; no BAO / nonlinear.",14,False,INK),
 ("• Idealized diagonal noise; no galaxy bias / RSD / survey window.",14,False,INK),
 ("• H₀ weak by design (h-unit shape degeneracy).",14,False,INK)])

# ---------- S11 TL;DR ----------
s=blank(); bg(s,NAVY)
tb(s,0.8,0.7,11.7,0.9,[("Take-home",34,True,WHITE)])
tb(s,0.8,1.9,11.7,4.4,[
 ("1.  We built a calibrated, amortized neural posterior estimator for ΛCDM parameters from the linear P(k).",18,True,ICE),
 ("",8,False,ICE),
 ("2.  Ωₘ, nₛ, Aₛ are well constrained and recovered; SBC + coverage confirm calibration.",18,True,ICE),
 ("",8,False,ICE),
 ("3.  H₀ is weakly constrained with a strong H₀–Ωₘ degeneracy (−0.96) — the correct shape-parameter physics.",18,True,ICE)])
tb(s,0.8,6.4,11.7,0.7,[("Maps the whole course: Monte Carlo · priors · PPC · SBC · likelihood-free · NPE.   Contact: [email]",14,False,GOLD)])

prs.save(OUTP); print("saved", OUTP, "slides:", len(prs.slides._sldIdLst))
